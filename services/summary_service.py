import aiofiles
from predibase import Predibase
from docx import Document
from pptx import Presentation
import PyPDF2
import os
from motor.motor_asyncio import AsyncIOMotorClient
from dotenv import load_dotenv

from bson.binary import Binary
import uuid

from bson import ObjectId
from bson.binary import UuidRepresentation


load_dotenv()

PREDIBASE_API_KEY = os.getenv("PREDIBASE_API_KEY")
MONGO_URI = os.getenv("MONGO_URI")

predibase_client = Predibase(api_token=PREDIBASE_API_KEY)

client = AsyncIOMotorClient(MONGO_URI)
db = client["filedb"]
collection = db["files"]


async def save_file(file, filename, file_extension):
    """
    Asynchronously saves a file to the specified storage folder.

    Args:
        file (aiofiles.os.RawIOBase): The file object to be saved.
        filename (str): The name of the file to be saved.
        file_extension (str): The extension of the file to be saved.

    Returns:
        str: The path of the saved file if successful, otherwise an error message.

    Raises:
        Exception: If an error occurs during the file saving process.
    """
    storage_folder = "storage/"
    if not os.path.exists(storage_folder):
        os.makedirs(storage_folder)

    file_path = os.path.join(storage_folder, filename)

    try:
        async with aiofiles.open(file_path, "wb") as out_file:
            content = await file.read()
            await out_file.write(content)

        return file_path

    except Exception as e:
        return "An error occurred"


def extract_text(file_path, file_extension):
    """
    Extracts text from a file based on its file extension.

    Args:
        file_path (str): The path of the file.
        file_extension (str): The extension of the file.

    Returns:
        str: The extracted text from the file.

    Raises:
        Exception: If an error occurs during the text extraction process.
    """
    text = ""
    try:
        if file_extension == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text
        elif file_extension == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
        elif file_extension == ".pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in range(len(reader.pages)):
                    page = reader.pages[page]
                    text += page.extract_text()
        return text
    except Exception as e:
        print(f"An error occurred: {e}")
        return "Error extracting text"


def get_summary(text):
    """
    Generates a summary of the given text using the Mistral-7b deployment of the predibase_client.

    Parameters:
        text (str): The text to be summarized.

    Returns:
        str: The generated summary of the text. If an error occurs during the summary generation,
        returns the string "Error generating summary".
    """
    try:
        lorax_client = predibase_client.deployments.client("mistral-7b")
        summary = str(lorax_client.generate(text, max_new_tokens=100).generated_text)
        return summary
    except Exception as e:
        print(f"An error occurred: {e}")
        return "Error generating summary"


def uuid_to_bson_binary(uuid_obj):
    """
    Converts a UUID object to a BSON binary representation.

    Args:
        uuid_obj (uuid.UUID): The UUID object to be converted.

    Returns:
        Binary: The BSON binary representation of the UUID object.

    Raises:
        TypeError: If the input is not a UUID object.

    This function takes a UUID object as input and converts it to a BSON binary object.
    BSON binary objects are used to store binary data in MongoDB.
    The UUID object is converted to a BSON binary object using the UuidRepresentation.STANDARD specification.
    """

    # Convert the UUID object to a BSON binary object using the UuidRepresentation.STANDARD specification.
    # The UuidRepresentation.STANDARD specification is used to ensure that the UUID object is converted to a BSON binary object in the standard format.
    return Binary.from_uuid(uuid_obj, UuidRepresentation.STANDARD)


def bson_binary_to_uuid(binary_obj):
    """
    Convert a BSON binary object to a UUID.

    This function takes a BSON binary object as input and converts it to a UUID object.
    The input can be of type Binary or ObjectId. If the input is of type Binary,
    it is converted to a UUID object using the bytes of the Binary object.
    If the input is of type ObjectId, it is converted to a string representation of the ObjectId.

    Args:
        binary_obj (Binary or ObjectId): The BSON binary object to convert.

    Returns:
        uuid.UUID or str: The converted UUID object or string representation of the ObjectId.

    Raises:
        TypeError: If the input object is not of type Binary or ObjectId.
    """
    if isinstance(binary_obj, Binary):
        # If the input is of type Binary, convert it to a UUID object using the bytes of the Binary object
        return uuid.UUID(bytes=binary_obj)
    elif isinstance(binary_obj, ObjectId):
        # If the input is of type ObjectId, convert it to a string representation of the ObjectId
        return str(binary_obj)
    else:
        # If the input is neither Binary nor ObjectId, raise a TypeError
        raise TypeError("Unsupported type for UUID conversion")
